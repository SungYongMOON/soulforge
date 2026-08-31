"""Bounded synthetic PPTX Workshop MCP used by the first physical file pilot."""

from __future__ import annotations

import argparse
import hashlib
import json
import os
from pathlib import Path
import re
import shutil

from mcp.server import MCPServer
from pptx import Presentation


REF = re.compile(r"^[A-Z][A-Z0-9-]{1,63}$")
CHECKPOINT = re.compile(r"^C[0-9]{4}$")
REVISION = re.compile(r"^R[0-9]{4}$")
DISPLAY_VERSION = re.compile(r"^V[0-9]+\.[0-9]+$")


def _require(value: str, pattern: re.Pattern[str], label: str) -> str:
    text = str(value or "").strip()
    if not pattern.fullmatch(text):
        raise ValueError(f"{label}_invalid")
    return text


def _root(env_name: str) -> Path:
    raw = os.environ.get(env_name, "").strip()
    if not raw:
        raise RuntimeError(f"{env_name}_missing")
    root = Path(raw).resolve()
    root.mkdir(parents=True, exist_ok=True)
    return root


def _inside(root: Path, *parts: str) -> Path:
    candidate = root.joinpath(*parts).resolve()
    try:
        candidate.relative_to(root)
    except ValueError as exc:
        raise ValueError("path_escape") from exc
    return candidate


def _bindings(project_ref: str, artifact_ref: str, job_ref: str | None = None) -> dict[str, Path | str]:
    project = _require(project_ref, REF, "project_ref")
    artifact = _require(artifact_ref, REF, "artifact_ref")
    job = _require(job_ref, REF, "job_ref") if job_ref is not None else None
    job_root = _root("PPT_JOB_ROOT")
    product_root = _root("PPT_WORK_PRODUCT_ROOT")
    return {
        "project": project,
        "artifact": artifact,
        "job": job,
        "job_dir": _inside(job_root, project, job) if job is not None else None,
        "artifact_dir": _inside(product_root, project, artifact),
    }


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for block in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest().upper()


def _slide_titles(path: Path) -> list[str]:
    presentation = Presentation(str(path))
    titles: list[str] = []
    for slide in presentation.slides:
        title = slide.shapes.title
        titles.append(title.text if title is not None else "")
    return titles


def _write_json_create_only(path: Path, payload: dict) -> None:
    if path.exists():
        raise FileExistsError(f"receipt_exists:{path.name}")
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True), encoding="utf-8")


def _require_job_binding(binding: dict[str, Path | str]) -> None:
    job_dir = binding["job_dir"]
    assert isinstance(job_dir, Path)
    request_path = job_dir / "REQUEST" / "request.json"
    if not request_path.is_file():
        raise FileNotFoundError("job_request_missing")
    try:
        request = json.loads(request_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise RuntimeError("job_request_invalid") from exc
    expected = {
        "project_ref": binding["project"],
        "artifact_ref": binding["artifact"],
        "job_ref": binding["job"],
    }
    if any(request.get(key) != value for key, value in expected.items()):
        raise PermissionError("job_binding_mismatch")


def initialize_synthetic_ppt(project_ref: str, artifact_ref: str, job_ref: str) -> dict:
    binding = _bindings(project_ref, artifact_ref, job_ref)
    job_dir = binding["job_dir"]
    artifact_dir = binding["artifact_dir"]
    assert isinstance(job_dir, Path) and isinstance(artifact_dir, Path)
    checkpoint_dir = artifact_dir / "CHECKPOINTS" / "C0000"
    checkpoint_file = checkpoint_dir / "pilot_deck.pptx"
    if checkpoint_file.exists():
        raise FileExistsError("artifact_exists")

    for name in ("REQUEST", "INPUT", "WORK", "OUTPUT", "CHECKPOINTS", "VALIDATION", "RECEIPT"):
        (job_dir / name).mkdir(parents=True, exist_ok=True)
    for name in ("WORKING", "CHECKPOINTS", "REVISIONS", "VALIDATION", "RECEIPTS"):
        (artifact_dir / name).mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    for number in range(1, 5):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"Baseline slide {number}"
    checkpoint_dir.mkdir(parents=True, exist_ok=False)
    presentation.save(str(checkpoint_file))
    working_file = artifact_dir / "WORKING" / "pilot_deck_CURRENT.pptx"
    shutil.copy2(checkpoint_file, working_file)
    job_checkpoint_dir = job_dir / "CHECKPOINTS" / "C0000"
    job_checkpoint_dir.mkdir(parents=True, exist_ok=False)
    shutil.copy2(checkpoint_file, job_checkpoint_dir / "pilot_deck.pptx")
    _write_json_create_only(
        job_dir / "REQUEST" / "request.json",
        {
            "schema_version": "soulforge.synthetic_ppt_job_request.v1",
            "project_ref": binding["project"],
            "artifact_ref": binding["artifact"],
            "job_ref": binding["job"],
            "kind": "synthetic_ppt_pilot",
            "authority_ceiling": "workshop_output_candidate_only",
            "real_project_payload": False,
            "human_acceptance": False,
            "backup_acceptance": False,
        },
    )
    payload = {
        "project_ref": binding["project"],
        "artifact_ref": binding["artifact"],
        "job_ref": binding["job"],
        "checkpoint_id": "C0000",
        "parent_checkpoint": None,
        "sha256": _sha256(checkpoint_file),
        "slide_titles": _slide_titles(checkpoint_file),
        "claim": "synthetic_working_checkpoint_only",
        "effect_count": 4,
    }
    _write_json_create_only(artifact_dir / "RECEIPTS" / "C0000.json", payload)
    return payload


def edit_slide_checkpoint(
    project_ref: str,
    artifact_ref: str,
    job_ref: str,
    parent_checkpoint: str,
    checkpoint_id: str,
    slide_number: int,
    replacement_text: str,
) -> dict:
    binding = _bindings(project_ref, artifact_ref, job_ref)
    _require_job_binding(binding)
    artifact_dir = binding["artifact_dir"]
    assert isinstance(artifact_dir, Path)
    parent = _require(parent_checkpoint, CHECKPOINT, "parent_checkpoint")
    checkpoint = _require(checkpoint_id, CHECKPOINT, "checkpoint_id")
    if not isinstance(slide_number, int) or not 1 <= slide_number <= 4:
        raise ValueError("slide_number_invalid")
    replacement = str(replacement_text or "").strip()
    if not replacement or len(replacement) > 200:
        raise ValueError("replacement_text_invalid")
    parent_file = artifact_dir / "CHECKPOINTS" / parent / "pilot_deck.pptx"
    if not parent_file.is_file():
        raise FileNotFoundError("parent_checkpoint_missing")
    target_dir = artifact_dir / "CHECKPOINTS" / checkpoint
    if target_dir.exists():
        raise FileExistsError("checkpoint_exists")
    target_dir.mkdir(parents=True, exist_ok=False)
    target_file = target_dir / "pilot_deck.pptx"

    presentation = Presentation(str(parent_file))
    slide = presentation.slides[slide_number - 1]
    if slide.shapes.title is None:
        raise RuntimeError("slide_title_missing")
    slide.shapes.title.text = replacement
    presentation.save(str(target_file))
    titles = _slide_titles(target_file)
    if titles[slide_number - 1] != replacement:
        raise RuntimeError("checkpoint_readback_mismatch")
    shutil.copy2(target_file, artifact_dir / "WORKING" / "pilot_deck_CURRENT.pptx")
    job_dir = binding["job_dir"]
    assert isinstance(job_dir, Path)
    job_checkpoint_dir = job_dir / "CHECKPOINTS" / checkpoint
    job_checkpoint_dir.mkdir(parents=True, exist_ok=False)
    shutil.copy2(target_file, job_checkpoint_dir / "pilot_deck.pptx")
    payload = {
        "project_ref": binding["project"],
        "artifact_ref": binding["artifact"],
        "job_ref": binding["job"],
        "checkpoint_id": checkpoint,
        "parent_checkpoint": parent,
        "changed_slide": slide_number,
        "sha256": _sha256(target_file),
        "slide_titles": titles,
        "claim": "synthetic_working_checkpoint_only",
        "effect_count": 3,
    }
    _write_json_create_only(artifact_dir / "RECEIPTS" / f"{checkpoint}.json", payload)
    return payload


def finalize_candidate_revision(
    project_ref: str,
    artifact_ref: str,
    job_ref: str,
    checkpoint_id: str,
    revision_id: str,
    display_version: str,
) -> dict:
    binding = _bindings(project_ref, artifact_ref, job_ref)
    _require_job_binding(binding)
    artifact_dir = binding["artifact_dir"]
    assert isinstance(artifact_dir, Path)
    checkpoint = _require(checkpoint_id, CHECKPOINT, "checkpoint_id")
    revision = _require(revision_id, REVISION, "revision_id")
    display = _require(display_version, DISPLAY_VERSION, "display_version")
    source = artifact_dir / "CHECKPOINTS" / checkpoint / "pilot_deck.pptx"
    if not source.is_file():
        raise FileNotFoundError("checkpoint_missing")
    revision_dir = artifact_dir / "REVISIONS" / revision
    if revision_dir.exists():
        raise FileExistsError("revision_exists")
    revision_dir.mkdir(parents=True, exist_ok=False)
    target = revision_dir / f"pilot_deck_{display}.pptx"
    shutil.copy2(source, target)
    source_hash = _sha256(source)
    target_hash = _sha256(target)
    if source_hash != target_hash:
        raise RuntimeError("revision_hash_mismatch")
    titles = _slide_titles(target)
    job_dir = binding["job_dir"]
    assert isinstance(job_dir, Path)
    job_output_dir = job_dir / "OUTPUT" / revision
    job_output_dir.mkdir(parents=True, exist_ok=False)
    job_output = job_output_dir / target.name
    shutil.copy2(target, job_output)
    if _sha256(job_output) != target_hash:
        raise RuntimeError("job_output_hash_mismatch")
    payload = {
        "project_ref": binding["project"],
        "artifact_ref": binding["artifact"],
        "job_ref": binding["job"],
        "source_checkpoint": checkpoint,
        "revision_id": revision,
        "display_version": display,
        "sha256": target_hash,
        "slide_titles": titles,
        "claim": "workshop_output_candidate_only",
        "human_acceptance": False,
        "backup_acceptance": False,
        "effect_count": 4,
    }
    _write_json_create_only(artifact_dir / "RECEIPTS" / f"{revision}.json", payload)
    _write_json_create_only(job_dir / "RECEIPT" / f"{revision}.json", payload)
    return payload


def get_ppt_checkpoint_state(project_ref: str, artifact_ref: str, checkpoint_id: str) -> dict:
    binding = _bindings(project_ref, artifact_ref)
    artifact_dir = binding["artifact_dir"]
    assert isinstance(artifact_dir, Path)
    checkpoint = _require(checkpoint_id, CHECKPOINT, "checkpoint_id")
    target = artifact_dir / "CHECKPOINTS" / checkpoint / "pilot_deck.pptx"
    if not target.is_file():
        raise FileNotFoundError("checkpoint_missing")
    return {
        "project_ref": binding["project"],
        "artifact_ref": binding["artifact"],
        "checkpoint_id": checkpoint,
        "sha256": _sha256(target),
        "slide_titles": _slide_titles(target),
        "effect_count": 0,
    }


def verify_ppt_revision(
    project_ref: str,
    artifact_ref: str,
    revision_id: str,
    expected_sha256: str,
) -> dict:
    binding = _bindings(project_ref, artifact_ref)
    artifact_dir = binding["artifact_dir"]
    assert isinstance(artifact_dir, Path)
    revision = _require(revision_id, REVISION, "revision_id")
    expected = str(expected_sha256 or "").strip().upper()
    if not re.fullmatch(r"[A-F0-9]{64}", expected):
        raise ValueError("expected_sha256_invalid")
    files = list((artifact_dir / "REVISIONS" / revision).glob("*.pptx"))
    if len(files) != 1:
        raise FileNotFoundError("revision_file_missing")
    actual = _sha256(files[0])
    failures = [] if actual == expected else ["SHA256_MISMATCH"]
    return {
        "verdict": "PASS" if not failures else "HOLD",
        "failures": failures,
        "project_ref": binding["project"],
        "artifact_ref": binding["artifact"],
        "revision_id": revision,
        "sha256": actual,
        "slide_titles": _slide_titles(files[0]),
        "authority_ceiling": "verified_completion_candidate",
        "effect_count": 0,
    }


mcp = MCPServer("ppt-workshop-pilot")
mcp.tool()(initialize_synthetic_ppt)
mcp.tool()(edit_slide_checkpoint)
mcp.tool()(finalize_candidate_revision)
mcp.tool()(get_ppt_checkpoint_state)
mcp.tool()(verify_ppt_revision)


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--transport", choices=("stdio", "http"), default="stdio")
    parser.add_argument("--port", type=int, default=18766)
    args = parser.parse_args()
    if args.transport == "http":
        mcp.run(
            transport="streamable-http",
            host="127.0.0.1",
            port=args.port,
            stateless_http=True,
        )
    else:
        mcp.run(transport="stdio")
